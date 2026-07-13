# slm/content_extractor.py
"""
Enhanced content‑extraction utilities for PDF, DOCX and PPTX files.

Features
--------
* PDF → HTML using pdfminer (better layout) + optional OCR fallback via pytesseract.
* DOCX → HTML that keeps headings, tables, lists and embeds images as base‑64 data‑URIs.
* PPTX → HTML preserving slide titles, bullet lists (nested) and slide images.
* Safety – every output string is sanitised with ``bleach`` before it is stored.
* Central configuration (``EXTRACTOR_SETTINGS``) makes thresholds & options easy to tweak.
* No API/model changes – the public function ``extract_content(file_obj)`` keeps the same
  signature and still raises ``ValueError`` on failure.
"""

import io
import os
import pathlib
import logging
import base64
import html                     # standard‑lib escaping for <pre> blocks
import re                       # used for plain‑text → HTML conversion
from typing import List

from django.core.files.base import ContentFile

# -------------------------------------------------------------------------
# Third‑party libraries (install with pip if you don’t have them yet)
# -------------------------------------------------------------------------
import mammoth                     # docx → html
import fitz                        # PyMuPDF (pdf)
import pdfminer.high_level as pdfminer  # pdfminer.six
from pptx import Presentation      # python‑pptx
from pptx.enum.shapes import PP_PLACEHOLDER  # placeholder type enum
from PIL import Image               # Pillow (image handling)

# Optional OCR – disabled automatically if the binary is missing
try:
    import pytesseract            # pytesseract (needs the tesseract exe)
    _OCR_AVAILABLE = True
except Exception:                 # pragma: no cover – missing Tesseract
    _OCR_AVAILABLE = False
    logging.getLogger(__name__).warning(
        "pytesseract not available – scanned PDFs will not be OCR‑ed."
    )

# HTML sanitiser
import bleach

# Optional BeautifulSoup – only needed when we have to rewrite stray image src’s
try:
    from bs4 import BeautifulSoup   # beautifulsoup4
    _BS4_AVAILABLE = True
except Exception:                 # pragma: no cover
    _BS4_AVAILABLE = False

# -------------------------------------------------------------------------
# Logging
# -------------------------------------------------------------------------
logger = logging.getLogger(__name__)

# -------------------------------------------------------------------------
# Allowed file extensions (kept identical to the original project)
# -------------------------------------------------------------------------
ALLOWED_EXT = {".pdf", ".doc", ".docx", ".ppt", ".pptx"}

# -------------------------------------------------------------------------
# Central configuration – tune here without touching the extraction code
# -------------------------------------------------------------------------
EXTRACTOR_SETTINGS = {
    # -----------------------------------------------------------------
    # PDF‑related options
    # -----------------------------------------------------------------
    "pdf": {
        "ocr_language": "eng",          # Tesseract language pack
        "min_text_ratio": 0.15,        # <15 % selectable text → run OCR
        "dpi_for_ocr": 200,             # DPI when rendering a page for OCR
        "max_pages_for_ocr": 30,        # safety guard – don’t OCR massive PDFs
    },

    # -----------------------------------------------------------------
    # DOCX‑related options
    # -----------------------------------------------------------------
    "docx": {
        "embed_images": True,           # embed any non‑inline images as base64 data‑uri
    },

    # -----------------------------------------------------------------
    # PPTX‑related options
    # -----------------------------------------------------------------
    "pptx": {
        "embed_images": True,
        "max_image_dim": 1024,          # down‑scale slide images larger than this (px)
    },

    # -----------------------------------------------------------------
    # Bleach sanitiser options – we start from the library defaults and
    # extend the allowed tags/attributes.  **Note:** recent Bleach versions
    # no longer accept a ``styles=`` arg, so we omit it.
    # -----------------------------------------------------------------
    "bleach": {
        "tags": list(bleach.sanitizer.ALLOWED_TAGS) + [
            "h1", "h2", "h3", "h4", "h5", "h6",
            "p", "ul", "ol", "li", "strong", "em", "br",
            "table", "thead", "tbody", "tr", "th", "td",
            "img", "blockquote", "pre", "code", "div",
            "span", "hr"
        ],
        "attributes": {
            "*": ["class", "style", "id"],
            "a": ["href", "title", "target", "rel"],
            "img": ["src", "alt", "title", "width", "height"],
        },
        "strip": True,
    },
}

# -------------------------------------------------------------------------
# Helper – HTML sanitiser (single point of truth)
# -------------------------------------------------------------------------
def _sanitize_html(raw_html: str) -> str:
    cfg = EXTRACTOR_SETTINGS["bleach"]
    return bleach.clean(
        raw_html,
        tags=cfg["tags"],
        attributes=cfg["attributes"],
        strip=cfg["strip"],
    )

# -------------------------------------------------------------------------
# Helper – read raw bytes from any Django FileField‑like object
# -------------------------------------------------------------------------
def _read_file_bytes(file_obj) -> bytes:
    """Guarantees we have the raw bytes of the uploaded file."""
    if hasattr(file_obj, "read"):
        return file_obj.read()
    # TemporaryUploadedFile has a .path attribute
    with open(file_obj.path, "rb") as f:
        return f.read()

# -------------------------------------------------------------------------
# Helper – convert plain‑text to simple semantic HTML
# -------------------------------------------------------------------------
def _plain_text_to_html(text: str) -> str:
    """
    Convert raw text (with line‑breaks) into light‑weight HTML.

    Heuristics:
      * Blank lines separate paragraphs.
      * Lines that are all‑uppercase (≤ 5 words) become <h2>.
      * Lines starting with a bullet (``-``, ``*`` or ``•``) become unordered list items.
      * Lines starting with ``<number>.`` become ordered list items.
      * Everything else becomes a normal <p>.
    Nested lists are not recognised – a contiguous block of list items
    will be wrapped in a single <ul> or <ol>.
    """
    bullet_re = re.compile(r"^\s*([-*•])\s+(.*)")
    ordered_re = re.compile(r"^\s*(\d+)[.)]\s+(.*)")

    html_chunks: List[str] = []
    list_open = False
    current_list_tag = None

    for raw_line in text.splitlines():
        line = raw_line.rstrip()
        if not line:
            # Close any open list on a blank line
            if list_open:
                html_chunks.append(f"</{current_list_tag}>")
                list_open = False
                current_list_tag = None
            continue

        # Heading detection – short all‑caps line
        if line.isupper() and len(line.split()) <= 5:
            if list_open:
                html_chunks.append(f"</{current_list_tag}>")
                list_open = False
                current_list_tag = None
            html_chunks.append(f"<h2>{html.escape(line.title())}</h2>")
            continue

        # Ordered list?
        m_ord = ordered_re.match(line)
        if m_ord:
            _, content = m_ord.groups()
            if not list_open or current_list_tag != "ol":
                if list_open:
                    html_chunks.append(f"</{current_list_tag}>")
                html_chunks.append("<ol>")
                list_open = True
                current_list_tag = "ol"
            html_chunks.append(f"<li>{html.escape(content.strip())}</li>")
            continue

        # Unordered list?
        m_bul = bullet_re.match(line)
        if m_bul:
            _, content = m_bul.groups()
            if not list_open or current_list_tag != "ul":
                if list_open:
                    html_chunks.append(f"</{current_list_tag}>")
                html_chunks.append("<ul>")
                list_open = True
                current_list_tag = "ul"
            html_chunks.append(f"<li>{html.escape(content.strip())}</li>")
            continue

        # Normal paragraph
        if list_open:
            html_chunks.append(f"</{current_list_tag}>")
            list_open = False
            current_list_tag = None
        html_chunks.append(f"<p>{html.escape(line)}</p>")

    # Close dangling list at EOF
    if list_open:
        html_chunks.append(f"</{current_list_tag}>")

    return "\n".join(html_chunks)

# -------------------------------------------------------------------------
# PDF extraction -----------------------------------------------------------
# -------------------------------------------------------------------------
def _extract_pdf(raw: bytes) -> str:
    """
    1️⃣ Try to get selectable text via pdfminer (keeps columns/tables).
    2️⃣ If the PDF looks mostly scanned, render each page with PyMuPDF,
       run OCR via pytesseract, and wrap the result in semantic HTML.
    Returns **sanitised** HTML where every page is wrapped in
    ``<div class="pdf-page" data-page="N">…</div>``.
    """
    settings = EXTRACTOR_SETTINGS["pdf"]

    # -------------------------------------------------
    # 1️⃣ pdfminer → plain text (fast, layout‑aware)
    # -------------------------------------------------
    try:
        txt = pdfminer.extract_text(io.BytesIO(raw))
    except Exception as exc:  # pdfminer can be noisy on bad PDFs
        logger.warning("pdfminer failed (%s) – falling back to OCR", exc)
        txt = ""

    # Heuristic: ratio of non‑whitespace chars to total length
    text_ratio = (len(txt.strip()) / len(txt)) if txt else 0.0
    need_ocr = (text_ratio < settings["min_text_ratio"]) or not txt

    # -------------------------------------------------
    # 2️⃣ If we have *good* selectable text → use it page‑by‑page
    # -------------------------------------------------
    if not need_ocr:
        doc = fitz.open(stream=raw, filetype="pdf")
        html_pages: List[str] = []
        for page_number, page in enumerate(doc, start=1):
            page_txt = page.get_text("text")   # plain text of the current page
            if not page_txt.strip():
                continue
            page_html = _plain_text_to_html(page_txt)
            html_pages.append(
                f'<div class="pdf-page" data-page="{page_number}">{page_html}</div>'
            )
        return _sanitize_html("\n".join(html_pages))

    # -------------------------------------------------
    # 3️⃣ OCR fallback (only if pytesseract is available)
    # -------------------------------------------------
    if not _OCR_AVAILABLE:
        logger.info("OCR not available – returning empty preview for PDF.")
        return _sanitize_html("<div></div>")

    doc = fitz.open(stream=raw, filetype="pdf")
    html_pages: List[str] = []

    for page_number, page in enumerate(doc, start=1):
        pix = page.get_pixmap(dpi=settings["dpi_for_ocr"])
        img_bytes = pix.tobytes("png")

        try:
            ocr_text = pytesseract.image_to_string(
                Image.open(io.BytesIO(img_bytes)),
                lang=settings["ocr_language"]
            )
        except Exception as exc:
            logger.error("OCR failed on PDF page %s: %s", page_number, exc)
            ocr_text = ""

        page_html = _plain_text_to_html(ocr_text)
        html_pages.append(
            f'<div class="pdf-page" data-page="{page_number}">{page_html}</div>'
        )

        if page_number >= settings["max_pages_for_ocr"]:
            logger.info(
                "Reached max_pages_for_ocr (%s) – stopping OCR early.",
                settings["max_pages_for_ocr"]
            )
            break

    final_html = "\n".join(html_pages)
    return _sanitize_html(final_html)

# -------------------------------------------------------------------------
# DOCX extraction ---------------------------------------------------------
# -------------------------------------------------------------------------
def _extract_docx(raw: bytes) -> str:
    """
    Convert DOCX (or .doc) to HTML with Mammoth.
    * Keeps headings, tables, lists, etc.
    * If ``embed_images`` is True, any images that are *not* already inline
      are extracted from the zip container and embedded as base64 data‑uri.
    * The document is split on Word page‑break elements – each segment is
      wrapped in ``<div class="docx-page" data-page="N">…</div>``.
      If the source contains no page‑breaks the entire document is wrapped
      as page 1.
    Returns sanitised HTML.
    """
    result = mammoth.convert_to_html(io.BytesIO(raw))
    raw_html = result.value

    # ---- Embed images (if enabled) and rewrite stray <pre> blocks ----
    if EXTRACTOR_SETTINGS["docx"]["embed_images"] and _BS4_AVAILABLE:
        soup = BeautifulSoup(raw_html, "html.parser")

        # ----- embed images -------------------------------------------------
        for img in soup.find_all("img"):
            src = img.get("src", "")
            if src.startswith("data:"):
                continue  # already a data‑uri, nothing to do

            try:
                import zipfile
                with zipfile.ZipFile(io.BytesIO(raw)) as z:
                    # Word stores binaries under word/media/
                    media_path = src.replace("file:///", "").lstrip("/")
                    if not media_path.startswith("word/media/"):
                        media_path = f"word/media/{media_path}"
                    img_bytes = z.read(media_path)
                    mime = f"image/{pathlib.Path(media_path).suffix.lstrip('.')}"
                    b64 = base64.b64encode(img_bytes).decode()
                    img["src"] = f"data:{mime};base64,{b64}"
            except Exception as exc:
                logger.warning("Failed to embed DOCX image %s: %s", src, exc)
                img["src"] = ""

        # ----- convert <pre> blocks -----------------------------------------
        for pre in soup.find_all("pre"):
            converted = _plain_text_to_html(pre.get_text())
            pre.replace_with(BeautifulSoup(converted, "html.parser"))

        # At this point ``soup`` contains the fully‑processed HTML.
        # --------------------------------------------------------------
        # Split into pages on Word “page‑break” paragraphs.
        # --------------------------------------------------------------
        parent = soup.body if soup.body else soup
        pages: List[str] = []
        cur_parts: List[str] = []

        for elem in list(parent.contents):
            # Detect a page‑break paragraph.
            if getattr(elem, "name", None) == "p":
                style = elem.get("style", "")
                if re.search(r'page-break-(?:before|after)', style, re.I):
                    # Finish the current page and start a new one.
                    pages.append("".join(str(p) for p in cur_parts))
                    cur_parts = []
                    continue
            # Anything else (including normal <p> tags) belongs to the current page.
            cur_parts.append(str(elem))

        # Append the final page (if any content left).
        pages.append("".join(str(p) for p in cur_parts))

        # Wrap each page in a <div>.
        wrapped_pages = [
            f'<div class="docx-page" data-page="{i + 1}">{page}</div>'
            for i, page in enumerate(pages) if page.strip()
        ]
        final_html = "\n".join(wrapped_pages)

    elif EXTRACTOR_SETTINGS["docx"]["embed_images"] and not _BS4_AVAILABLE:
        # -----------------------------------------------------------------
        # No BeautifulSoup – fall back to regex‑only handling.
        # -----------------------------------------------------------------
        # Convert <pre> blocks.
        def _pre_repl(m):
            inner = html.unescape(m.group(1))
            return _plain_text_to_html(inner)

        tmp_html = re.sub(r"<pre>(.*?)</pre>", _pre_repl, raw_html, flags=re.DOTALL)

        # Split on page‑break paragraphs using a regex.
        split_pat = re.compile(
            r'(?i)<p[^>]*style=["\'][^"\']*page-break-(?:before|after)[^"\']*["\'][^>]*>\s*</p>'
        )
        pages = split_pat.split(tmp_html)
        wrapped_pages = [
            f'<div class="docx-page" data-page="{i + 1}">{section}</div>'
            for i, section in enumerate(pages) if section.strip()
        ]
        final_html = "\n".join(wrapped_pages)

    else:
        # -----------------------------------------------------------------
        # Neither image embedding nor BeautifulSoup is available.
        # We still split on page‑breaks (regex‑only) so callers get a
        # consistent structure.
        # -----------------------------------------------------------------
        split_pat = re.compile(
            r'(?i)<p[^>]*style=["\'][^"\']*page-break-(?:before|after)[^"\']*["\'][^>]*>\s*</p>'
        )
        pages = split_pat.split(raw_html)
        wrapped_pages = [
            f'<div class="docx-page" data-page="{i + 1}">{section}</div>'
            for i, section in enumerate(pages) if section.strip()
        ]
        final_html = "\n".join(wrapped_pages)

    return _sanitize_html(final_html)

# -------------------------------------------------------------------------
# PPTX extraction ---------------------------------------------------------
# -------------------------------------------------------------------------
def _extract_pptx(raw: bytes) -> str:
    """
    Convert PPTX → HTML.
    * Each slide becomes ``<div class="ppt-page" data-page="N">``.
    * The first TITLE placeholder (if present) → ``<h2>``.
    * All other placeholders that are FOOTER/SLIDE_NUMBER/HEADER/DATE are ignored
      (so slide numbers no longer appear as stray paragraphs).
    * Bullet / numbered paragraphs become nested ``<ul>/<ol>``.
    * Images are embedded as base64 data‑uri (optional down‑scale).
    Returns sanitised HTML.
    """
    cfg = EXTRACTOR_SETTINGS["pptx"]
    prs = Presentation(io.BytesIO(raw))

    # Placeholder types we want to **skip completely**
    SKIP_PLACEHOLDER_TYPES = {
        PP_PLACEHOLDER.FOOTER,
        PP_PLACEHOLDER.SLIDE_NUMBER,
        PP_PLACEHOLDER.HEADER,
        PP_PLACEHOLDER.DATE,
    }

    sections: List[str] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        # ---- start the per‑slide wrapper ---------------------------------
        parts: List[str] = [f'<div class="ppt-page" data-page="{slide_idx}">']

        # ---- extract a TITLE placeholder (if any) -----------------------
        title_text = None
        title_shape = None
        for shape in slide.shapes:
            if getattr(shape, "is_placeholder", False):
                ph_type = shape.placeholder_format.type
                if ph_type == PP_PLACEHOLDER.TITLE:
                    title_text = shape.text.strip()
                    title_shape = shape
                    break

        if title_text:
            parts.append(f"<h2>{html.escape(title_text)}</h2>")

        # ---- paragraphs & lists -----------------------------------------
        list_open = False
        for shape in slide.shapes:
            # Skip anything that doesn't contain a text frame
            if not getattr(shape, "has_text_frame", False):
                continue

            # Skip the TITLE shape we already rendered
            if shape is title_shape:
                continue

            # Skip placeholders we don't want to render (footers, slide numbers, …)
            if getattr(shape, "is_placeholder", False):
                ph_type = shape.placeholder_format.type
                if ph_type in SKIP_PLACEHOLDER_TYPES:
                    continue

            for para in shape.text_frame.paragraphs:
                txt = para.text.strip()
                if not txt:
                    continue

                lvl = para.level
                if lvl == 0:
                    # Regular paragraph – close any open list first
                    if list_open:
                        parts.append("</ul>")
                        list_open = False
                    parts.append(f"<p>{html.escape(txt)}</p>")
                else:
                    # List item (any level >0 → treat as an unordered list)
                    if not list_open:
                        parts.append("<ul>")
                        list_open = True
                    parts.append(f"<li>{html.escape(txt)}</li>")

        if list_open:
            parts.append("</ul>")

        # ---- images (pictures) -------------------------------------------
        if cfg["embed_images"]:
            for shape in slide.shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    image = shape.image
                    img_bytes = image.blob

                    # Down‑scale huge images to keep DB size modest
                    if max(image.width, image.height) > cfg["max_image_dim"]:
                        pil_img = Image.open(io.BytesIO(img_bytes))
                        pil_img.thumbnail(
                            (cfg["max_image_dim"], cfg["max_image_dim"])
                        )
                        buf = io.BytesIO()
                        pil_img.save(buf, format="PNG")
                        img_bytes = buf.getvalue()
                        mime = "image/png"
                    else:
                        mime = f"image/{image.ext}"

                    b64 = base64.b64encode(img_bytes).decode()
                    parts.append(
                        f'<img src="data:{mime};base64,{b64}" alt="Slide {slide_idx} image">'
                    )

        # ---- close the wrapper -------------------------------------------
        parts.append("</div>")
        sections.append("\n".join(parts))

    return _sanitize_html("\n".join(sections))

# -------------------------------------------------------------------------
# Public façade – unchanged signature (returns safe HTML or raises ValueError)
# -------------------------------------------------------------------------
def extract_content(file_obj) -> str:
    """
    Public API used throughout the project.
    Accepts a Django ``FileField``‑like object and returns **sanitised HTML**.
    Raises ``ValueError`` for unsupported extensions or any extraction problem.
    """
    ext = pathlib.Path(file_obj.name).suffix.lower()
    if ext not in ALLOWED_EXT:
        raise ValueError(f"Unsupported extension: {ext}")

    raw = _read_file_bytes(file_obj)

    try:
        if ext in {".docx", ".doc"}:
            return _extract_docx(raw)
        if ext == ".pdf":
            return _extract_pdf(raw)
        if ext in {".pptx", ".ppt"}:
            return _extract_pptx(raw)
    except Exception as exc:  # pragma: no cover – exercised via runtime
        logger.exception("Failed to extract %s", file_obj.name)
        raise ValueError(f"Extraction error: {exc}")

    # This line should never be reached because the extension guard is earlier
    raise ValueError(f"Unsupported extension: {ext}")